import os
import re
import json
import subprocess
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Ensure environment has needed libraries: pip install PyPDF2 python-docx python-pptx pandas
# Using absolute path for rclone
RCLONE_BIN = "/opt/homebrew/bin/rclone"
TEMP_DIR = "/tmp/twister_analysis"

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for i, page in enumerate(reader.pages):
            if i > 2: break # Only read first 3 pages
            text += page.extract_text() + " "
        return text.strip()
    except Exception as e:
        return f"Error: {e}"

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = ""
        for i, p in enumerate(doc.paragraphs):
            if i > 50: break # Limit
            text += p.text + " "
        return text.strip()
    except Exception as e:
        return f"Error: {e}"

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            if i > 5: break # Only read first 5 slides
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text.strip()
    except Exception as e:
        return f"Error: {e}"

def generate_name_from_text(text, orig_ext):
    if not text or len(text.strip()) < 5:
        return None
        
    # Clean text: remove newlines, multiple spaces, weird chars
    clean_text = re.sub(r'[\r\n\t]', ' ', text)
    clean_text = re.sub(r'[^\w\s\-\.]', '', clean_text)
    clean_text = re.sub(r'\s+', ' ', clean_text).strip()
    
    words = clean_text.split()
    
    # Very basic heuristic: Look for a title-like sequence
    # Often the first few words are the title/name
    meaningful_words = []
    for w in words:
        if len(w) > 2 and w.lower() not in ['the', 'and', 'for', 'with', 'this', 'that', 'from']:
            meaningful_words.append(w.capitalize())
        if len(meaningful_words) >= 6:
            break
            
    if not meaningful_words:
        return None
        
    new_base = " ".join(meaningful_words)
    return f"{new_base}{orig_ext}"

def main():
    print("Finding numeric files in Twister...")
    # List all files in Twister
    ls_cmd = [RCLONE_BIN, "lsjson", "-R", "gdrive:1.Organize/Twister"]
    ls_output = subprocess.check_output(ls_cmd)
    files = json.loads(ls_output)
    
    # Filter for numeric files that are docs/pdfs/ppts
    target_exts = ['.pdf', '.docx', '.pptx']
    numeric_pattern = re.compile(r'^(\d+)\.[a-zA-Z0-9]+$')
    
    target_files = []
    for f in files:
        if f.get('IsDir'): continue
        name = f.get('Name', '')
        ext = os.path.splitext(name)[1].lower()
        if ext in target_exts and numeric_pattern.match(name):
            target_files.append(f)
            
    print(f"Found {len(target_files)} numeric documents to analyze.")
    
    results = []
    
    for idx, f in enumerate(target_files):
        path = f["Path"]
        name = f["Name"]
        file_id = f["ID"]
        ext = os.path.splitext(name)[1].lower()
        
        remote_path = f"gdrive:1.Organize/Twister/{path}"
        local_path = os.path.join(TEMP_DIR, f"temp_{idx}{ext}")
        
        print(f"[{idx+1}/{len(target_files)}] Analyzing: {path}")
        
        try:
            # Download
            subprocess.run([RCLONE_BIN, "copyto", remote_path, local_path], check=True, capture_output=True)
            
            # Extract Text
            text = ""
            if ext == '.pdf':
                text = extract_text_from_pdf(local_path)
            elif ext == '.docx':
                text = extract_text_from_docx(local_path)
            elif ext == '.pptx':
                text = extract_text_from_pptx(local_path)
                
            # Generate Name
            new_name = generate_name_from_text(text, ext)
            
            if new_name:
                print(f"  -> Proposed name: {new_name}")
                
                # Perform the rename via moveid
                dir_path = os.path.dirname(f"1.Organize/Twister/{path}")
                new_remote_path = f"gdrive:{dir_path}/{new_name}"
                
                rename_cmd = [RCLONE_BIN, "backend", "moveid", "gdrive:", file_id, f"{dir_path}/{new_name}"]
                rename_res = subprocess.run(rename_cmd, capture_output=True, text=True)
                
                if rename_res.returncode == 0:
                    print("  -> Successfully renamed.")
                    status = "Renamed"
                else:
                    print(f"  -> Rename failed: {rename_res.stderr.strip()}")
                    status = f"Rename Failed: {rename_res.stderr.strip()}"
            else:
                print("  -> Could not generate meaningful name.")
                new_name = ""
                status = "No meaningful text"
                
            results.append({
                "OriginalPath": path,
                "OriginalName": name,
                "ProposedName": new_name,
                "Ext": ext,
                "Status": status,
                "Snippet": text[:200] if text else ""
            })
            
            # Cleanup
            if os.path.exists(local_path):
                os.remove(local_path)
                
        except Exception as e:
            print(f"  -> Error processing: {e}")
            results.append({
                "OriginalPath": path,
                "OriginalName": name,
                "ProposedName": "",
                "Ext": ext,
                "Status": f"Error: {str(e)}",
                "Snippet": ""
            })

    # Save report
    report_df = pd.DataFrame(results)
    report_path = "/Volumes/Fulcrum/Develop/gdrive-control/reports/twister_rename_report.csv"
    report_df.to_csv(report_path, index=False)
    print(f"\nAnalysis and renaming complete. Report saved to {report_path}")

if __name__ == "__main__":
    main()
