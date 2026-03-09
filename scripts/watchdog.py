import os
import re
import json
import subprocess
import pandas as pd
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from datetime import datetime

# Setup OpenAI
api_key = os.environ.get("OPENAI_API_KEY")
if not api_key:
    try:
        with open("/Volumes/Fulcrum/Develop/OHLC/.env", "r") as f:
            for line in f:
                if line.startswith("OPENAI_API_KEY="):
                    api_key = line.strip().split("=", 1)[1]
                    break
    except Exception as e:
        print(f"Could not load API key: {e}")

client = OpenAI(api_key=api_key)

RCLONE_BIN = "/opt/homebrew/bin/rclone"
TEMP_DIR = "/tmp/auto_renamer"
os.makedirs(TEMP_DIR, exist_ok=True)
LOG_PATH = "/Volumes/Fulcrum/Develop/gdrive-control/reports/auto_rename_log.csv"

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for i, page in enumerate(reader.pages):
            if i > 2: break
            text += page.extract_text() + " "
        return text.strip()
    except Exception:
        return ""

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = ""
        for i, p in enumerate(doc.paragraphs):
            if i > 50: break
            text += p.text + " "
        return text.strip()
    except Exception:
        return ""

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            if i > 5: break
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text.strip()
    except Exception:
        return ""

def generate_name_with_llm(original_name, folder_path, text_snippet):
    prompt = f"""
You are an expert file organizer. Your job is to rename a file based on its content and context.
The standard naming convention is: [Entity/Project] - [Core Subject] - [Date/Version] (if applicable).
Rules:
- Use spaces and hyphens for readability.
- Do NOT use camelCase or snake_case.
- Strip out random UUIDs, hashes, and generic prefixes like "Copy of", "Untitled", "Screen Shot".
- If a date is known, put it at the end.
- Return ONLY the new filename, including the original extension. Do not include any other text or explanation.
- If no text snippet is provided, use the folder path context to guess a reasonable name. Do NOT output literal brackets like '[Original Identifier]'.

Original Filename: {original_name}
Folder Path Context: {folder_path}
Text Snippet from File (if any):
{text_snippet[:1500]}
"""
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that renames files."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=60
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"LLM Error: {e}")
        return None

def is_odd_name(filename):
    name, ext = os.path.splitext(filename)
    name = name.strip()
    
    uuid_pattern = re.compile(r'^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$', re.IGNORECASE)
    hash_pattern = re.compile(r'^[0-9a-f]{20,}$', re.IGNORECASE)
    numeric_pattern = re.compile(r'^\d+$')
    generic_pattern = re.compile(r'^(untitled|copy of|document\d*|presentation\d*|book\d*|image\d*|screenshot|screen shot)', re.IGNORECASE)
    weird_chars_pattern = re.compile(r'^[^a-zA-Z0-9\s]+$')
    
    if numeric_pattern.match(name): return True
    if uuid_pattern.match(name) or hash_pattern.match(name): return True
    if len(name) <= 2: return True
    if generic_pattern.match(name): return True
    if weird_chars_pattern.match(name): return True
    if re.match(r'^[\~\!\@\#\$\%\^\&\*\(\)\_\+\[\]\{\}\;\:\,\>\<\?\|]', name): return True
    
    return False

def main():
    print(f"[{datetime.now().isoformat()}] Starting Naming Watchdog...")
    
    # Use rclone to find files modified in the last 24 hours
    cmd = [RCLONE_BIN, "lsjson", "-R", "--max-age", "24h", "--files-only", "gdrive:"]
    print("Scanning Drive for recent files...")
    result = subprocess.run(cmd, capture_output=True, text=True)
    
    if result.returncode != 0:
        print(f"Error scanning drive: {result.stderr}")
        return
        
    files = json.loads(result.stdout)
    print(f"Found {len(files)} files modified in the last 24 hours.")
    
    odd_files = []
    for f in files:
        if is_odd_name(f['Name']):
            odd_files.append(f)
            
    print(f"Identified {len(odd_files)} files with odd names.")
    
    if not odd_files:
        print("No action needed. Exiting.")
        return
        
    # Load existing log to avoid reprocessing if it failed halfway
    processed_ids = set()
    if os.path.exists(LOG_PATH):
        df_log = pd.read_csv(LOG_PATH)
        processed_ids = set(df_log['ID'].tolist())
        
    log_file = open(LOG_PATH, "a")
    if not os.path.exists(LOG_PATH) or os.path.getsize(LOG_PATH) == 0:
        log_file.write("ID,OriginalPath,OriginalName,NewName,Status,Snippet\n")
        
    text_exts = ['.pdf', '.docx', '.pptx']
    
    for idx, f in enumerate(odd_files):
        file_id = f.get('ID')
        if not file_id or file_id in processed_ids:
            continue
            
        path = f['Path']
        if "SupabaseBackups" in path:
            print(f"[{idx+1}/{len(odd_files)}] Skipping SupabaseBackups folder: {f['Name']}")
            continue
            
        name = f['Name']
        ext = os.path.splitext(name)[1].lower()
        
        print(f"[{idx+1}/{len(odd_files)}] Processing: {name}")
        
        text = ""
        if ext in text_exts:
            remote_path = f"gdrive:{path}"
            local_path = os.path.join(TEMP_DIR, f"temp_watchdog_{file_id}{ext}")
            
            try:
                subprocess.run([RCLONE_BIN, "copyto", remote_path, local_path], check=True, capture_output=True)
                if ext == '.pdf':
                    text = extract_text_from_pdf(local_path)
                elif ext == '.docx':
                    text = extract_text_from_docx(local_path)
                elif ext == '.pptx':
                    text = extract_text_from_pptx(local_path)
            except Exception as e:
                print(f"  -> Error downloading/extracting: {e}")
            finally:
                if os.path.exists(local_path):
                    os.remove(local_path)
                    
        folder_path = os.path.dirname(path)
        new_name = generate_name_with_llm(name, folder_path, text)
        
        if not new_name or new_name == name:
            print("  -> Could not generate a new name or name is the same.")
            safe_text = text[:100].replace('\n', ' ').replace(',', ' ') if text else ""
            log_file.write(f"{file_id},{path},{name},{name},Skipped,{safe_text}\n")
            log_file.flush()
            continue
            
        new_name = new_name.replace("/", "-").replace(":", "-").replace('"', '').replace("'", "")
        print(f"  -> Proposed Name: {new_name}")
        
        dest_path = f"{folder_path}/{new_name}" if folder_path else new_name
        
        try:
            cmd_move = [RCLONE_BIN, "backend", "moveid", "gdrive:", file_id, dest_path]
            res_move = subprocess.run(cmd_move, capture_output=True, text=True)
            
            if res_move.returncode == 0:
                print("  -> Success")
                safe_text = text[:100].replace('\n', ' ').replace(',', ' ') if text else ""
                log_file.write(f"{file_id},{path},{name},{new_name},Success,{safe_text}\n")
            else:
                print(f"  -> Rclone Error: {res_move.stderr.strip()}")
                safe_text = text[:100].replace('\n', ' ').replace(',', ' ') if text else ""
                log_file.write(f"{file_id},{path},{name},{new_name},Error,{safe_text}\n")
        except Exception as e:
            print(f"  -> Execution Error: {e}")
            
        log_file.flush()
        
    log_file.close()
    print("Watchdog run complete.")

if __name__ == "__main__":
    main()
