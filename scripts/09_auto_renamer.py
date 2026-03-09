import os
import re
import json
import subprocess
import pandas as pd
import glob
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import time

# Setup OpenAI
api_key = os.environ.get("OPENAI_API_KEY")
if not api_key:
    try:
        with open("/Volumes/Fulcrum/Develop/OHLC/.env", "r") as f:
            for line in f:
                if line.startswith("OPENAI_API_KEY="):
                    api_key = line.strip().split("=", 1)[1]
                    break
    except Exception as e:
        print(f"Could not load API key: {e}")

client = OpenAI(api_key=api_key)

RCLONE_BIN = "/opt/homebrew/bin/rclone"
TEMP_DIR = "/tmp/auto_renamer"
os.makedirs(TEMP_DIR, exist_ok=True)

def extract_text_from_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for i, page in enumerate(reader.pages):
            if i > 2: break
            text += page.extract_text() + " "
        return text.strip()
    except Exception:
        return ""

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = ""
        for i, p in enumerate(doc.paragraphs):
            if i > 50: break
            text += p.text + " "
        return text.strip()
    except Exception:
        return ""

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            if i > 5: break
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text.strip()
    except Exception:
        return ""

def generate_name_with_llm(original_name, folder_path, text_snippet):
    prompt = f"""
You are an expert file organizer. Your job is to rename a file based on its content and context.
The standard naming convention is: [Entity/Project] - [Core Subject] - [Date/Version] (if applicable).
Rules:
- Use spaces and hyphens for readability.
- Do NOT use camelCase or snake_case.
- Strip out random UUIDs, hashes, and generic prefixes like "Copy of" or "Untitled".
- If a date is known, put it at the end.
- Return ONLY the new filename, including the original extension. Do not include any other text or explanation.
- If no text snippet is provided, use the folder path context to guess a reasonable name (e.g., if in 'Softvision/Design', name it 'Softvision - Design - Document.ext'). Do NOT output literal brackets like '[Original Identifier]'.

Original Filename: {original_name}
Folder Path Context: {folder_path}
Text Snippet from File (if any):
{text_snippet[:1500]}
"""
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that renames files."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=60
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"LLM Error: {e}")
        return None

def main():
    report_path = "/Volumes/Fulcrum/Develop/gdrive-control/reports/odd_filenames_report.csv"
    if not os.path.exists(report_path):
        print("Odd filenames report not found.")
        return

    df_odd = pd.read_csv(report_path)
    
    # Load migration map to get IDs
    runs = sorted(glob.glob("/Volumes/Fulcrum/Develop/gdrive-control/runs/run_*"))
    latest_run = runs[-1]
    map_path = f"{latest_run}/07-migration/migration_map_resolved.csv"
    df_map = pd.read_csv(map_path)
    
    # Merge to get ID
    df = pd.merge(df_odd, df_map[['Path', 'ID']], on='Path', how='left')
    
    # Filter out files that already have no ID
    df = df.dropna(subset=['ID'])
    
    log_path = "/Volumes/Fulcrum/Develop/gdrive-control/reports/auto_rename_log.csv"
    
    # Load existing log if any to resume
    processed_ids = set()
    if os.path.exists(log_path):
        with open(log_path, 'r', encoding='utf-8', errors='replace') as f:
            for line in f:
                if line.strip() and not line.startswith('ID,'):
                    parts = line.split(',')
                    if parts:
                        processed_ids.add(parts[0])
    
    df_to_process = df[~df['ID'].isin(processed_ids)]
    
    print(f"Found {len(df_to_process)} files left to rename out of {len(df)} total odd files.")
    
    log_file = open(log_path, "a")
    if not os.path.exists(log_path) or os.path.getsize(log_path) == 0:
        log_file.write("ID,OriginalPath,OriginalName,NewName,Status,Snippet\n")
    
    text_exts = ['.pdf', '.docx', '.pptx']
    
    for idx, row in df_to_process.iterrows():
        path = row['Path']
        name = row['Name']
        ext = row['Ext']
        file_id = row['ID']
        
        if "SupabaseBackups" in path:
            print(f"[{idx+1}/{len(df_to_process)}] Skipping SupabaseBackups folder: {name}")
            continue
            
        print(f"[{idx+1}/{len(df_to_process)}] Processing: {name}")
        
        text = ""
        if ext in text_exts:
            remote_path = f"gdrive:{path}"
            local_path = os.path.join(TEMP_DIR, f"temp_{file_id}{ext}")
            
            try:
                subprocess.run([RCLONE_BIN, "copyto", remote_path, local_path], check=True, capture_output=True)
                if ext == '.pdf':
                    text = extract_text_from_pdf(local_path)
                elif ext == '.docx':
                    text = extract_text_from_docx(local_path)
                elif ext == '.pptx':
                    text = extract_text_from_pptx(local_path)
            except Exception as e:
                print(f"  -> Error downloading/extracting: {e}")
            finally:
                if os.path.exists(local_path):
                    os.remove(local_path)
        
        folder_path = os.path.dirname(path)
        new_name = generate_name_with_llm(name, folder_path, text)
        
        if not new_name or new_name == name:
            print("  -> Could not generate a new name or name is the same.")
            safe_text = text[:100].replace('\n', ' ').replace(',', ' ') if text else ""
            log_file.write(f"{file_id},{path},{name},{name},Skipped,{safe_text}\n")
            log_file.flush()
            continue
            
        # Clean new name just in case
        new_name = new_name.replace("/", "-").replace(":", "-").replace('"', '').replace("'", "")
        
        print(f"  -> Proposed Name: {new_name}")
        
        dest_path = f"{folder_path}/{new_name}" if folder_path else new_name
        
        # Rename using rclone backend moveid
        try:
            cmd = [RCLONE_BIN, "backend", "moveid", "gdrive:", file_id, dest_path]
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode == 0:
                print("  -> Success")
                safe_text = text[:100].replace('\n', ' ').replace(',', ' ') if text else ""
                log_file.write(f"{file_id},{path},{name},{new_name},Success,{safe_text}\n")
            else:
                print(f"  -> Rclone Error: {result.stderr.strip()}")
                safe_text = text[:100].replace('\n', ' ').replace(',', ' ') if text else ""
                log_file.write(f"{file_id},{path},{name},{new_name},Error,{safe_text}\n")
        except Exception as e:
            print(f"  -> Execution Error: {e}")
            
        log_file.flush()

    log_file.close()
    print("Done.")

if __name__ == "__main__":
    main()
