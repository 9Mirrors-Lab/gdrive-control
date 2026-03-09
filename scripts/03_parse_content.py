import os
import glob
import pandas as pd
from pathlib import Path
import json

# Parsing libraries
import PyPDF2
import docx
import pptx
import openpyxl
import csv

RUN_ID = "run_20260306_151014"
BASE_DIR = f"/Volumes/Fulcrum/Develop/gdrive-control/runs/{RUN_ID}"
STAGING_DIR = f"{BASE_DIR}/02-content/staging/batch-001"
OUTPUT_FILE = f"{BASE_DIR}/02-content/content_index.csv"
ERROR_LOG = f"{BASE_DIR}/02-content/extraction_failures.csv"

def extract_pdf(filepath):
    text = ""
    with open(filepath, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for i in range(min(3, len(reader.pages))):
            text += reader.pages[i].extract_text() + " "
    return text[:1000].strip()

def extract_docx(filepath):
    doc = docx.Document(filepath)
    text = " ".join([p.text for p in doc.paragraphs[:20]])
    return text[:1000].strip()

def extract_pptx(filepath):
    prs = pptx.Presentation(filepath)
    text = ""
    for slide in prs.slides[:3]:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text[:1000].strip()

def extract_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, data_only=True)
    text = ""
    sheet = wb.active
    for row in sheet.iter_rows(min_row=1, max_row=20, values_only=True):
        text += " ".join([str(c) for c in row if c]) + " "
    return text[:1000].strip()

def extract_csv(filepath):
    text = ""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i > 20: break
            text += " ".join(row) + " "
    return text[:1000].strip()

results = []
failures = []

for root, _, files in os.walk(STAGING_DIR):
    for file in files:
        filepath = os.path.join(root, file)
        rel_path = os.path.relpath(filepath, STAGING_DIR)
        ext = Path(file).suffix.lower()
        
        content = ""
        try:
            if ext == '.pdf':
                content = extract_pdf(filepath)
            elif ext == '.docx':
                content = extract_docx(filepath)
            elif ext == '.pptx':
                content = extract_pptx(filepath)
            elif ext == '.xlsx':
                content = extract_xlsx(filepath)
            elif ext == '.csv':
                content = extract_csv(filepath)
            elif ext in ['.txt', '.md', '.json']:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(1000)
            
            results.append({
                "Path": rel_path,
                "Extension": ext,
                "Snippet": content.replace("\n", " ")[:500]
            })
        except Exception as e:
            failures.append({
                "Path": rel_path,
                "Extension": ext,
                "Error": str(e)
            })

df = pd.DataFrame(results)
df.to_csv(OUTPUT_FILE, index=False)

df_err = pd.DataFrame(failures)
if not df_err.empty:
    df_err.to_csv(ERROR_LOG, index=False)

print(f"Extraction complete. Successfully parsed: {len(results)}, Failed: {len(failures)}")
